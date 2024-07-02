# -*- coding: utf-8 -*-
"""
Created on Sat Jun 29 17:15:30 2024

@author: piret
"""

# galora.py

"""
Created on Mon Jun 17 16:01:34 2024

Author: piret
"""

import json
import os
import logging
import argparse
from datetime import datetime
import fitz  # PyMuPDF
from pptx import Presentation
from moviepy.editor import VideoFileClip
import speech_recognition as sr
from pydub import AudioSegment
import pandas as pd
import csv
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
from pydub.silence import split_on_silence
from pytube import YouTube
from docx import Document
import zipfile
import re
from difflib import SequenceMatcher
import requests
import shutil
import vlc
import xml.etree.ElementTree as ET
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
import subprocess

# Global variable for language
lang = {}

# Configure logging
log_dir = "log"
temp_dir = "temp"
if not os.path.exists(log_dir):
    os.makedirs(log_dir)
if not os.path.exists(temp_dir):
    os.makedirs(temp_dir)

# 1. configure_logger: Configures the logger for the specified module
def configure_logger(module_name):
    """Configures the logger for the specified module."""
    timestamp = datetime.now().strftime("%Y%m%d")
    log_dir = "log"
    if not os.path.exists(log_dir):
        os.makedirs(log_dir)
    log_file = os.path.join(log_dir, f'{module_name}_{timestamp}.log')
    try:
        logging.basicConfig(
            filename=log_file,
            level=logging.DEBUG,
            format='%(asctime)s - %(levelname)s - %(message)s'
        )
        logging.info("Logger has been configured successfully.")
        print(f"Logger configured: {log_file}")
    except Exception as e:
        print(f"Failed to configure logger: {e}")
    return log_file

# 2. log_message: Logs messages using translation keys
def log_message(key, level="info", *args):
    """Logs messages using translation keys."""
    try:
        message = lang.get(key, key).format(*args)
    except KeyError:
        message = f"Logging key error: {key} with args {args}"

    if level == "info":
        logging.info(message)
    elif level == "warning":
        logging.warning(message)
    elif level == "error":
        logging.error(message)
    elif level == "debug":
        logging.debug(message)
    print(f"Log message: {message}")

# 3. load_config: Loads the configuration for the specified module
def load_config(module_name, config_path='config.json'):
    """Loads the configuration for the specified module."""
    try:
        with open(config_path, 'r') as config_file:
            configs = json.load(config_file)
            for config in configs:
                if config["module"] == module_name:
                    log_message('config_loaded', 'info', module_name)
                    return config
        log_message('config_module_not_found', 'warning', module_name)
    except json.JSONDecodeError as e:
        log_message('json_decode_error', 'error', str(e))
    except Exception as e:
        log_message('config_load_failed', 'error', str(e))
    return None

# 4. load_translations: Loads the translation file based on the language code
def load_translations(language_code):
    """Loads the translation file based on the language code."""
    try:
        with open(f'language/cli_{language_code}.json', 'r', encoding='utf-8') as lang_file:
            return json.load(lang_file)
    except FileNotFoundError:
        log_message('language_file_not_found', 'error', language_code)
        return {}
    except json.JSONDecodeError as e:
        log_message('json_decode_error_language', 'error', str(e))
        return {}
    except Exception as e:
        log_message('error_loading_language_file', 'error', str(e))
        return {}

# 5. remove_headers_footers: Removes headers and footers from text
def remove_headers_footers(text):
    """Removes headers and footers from text."""
    lines = text.split('\n')
    if len(lines) > 3:
        return '\n'.join(lines[1:-1])
    return text

# 6. handle_text_file: Processes text files
def handle_text_file(file_path):
    """Processes text files."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            text = file.read()
        log_message('text_file_processed', 'info', file_path)
        return remove_headers_footers(text), file_path
    except Exception as e:
        log_message('error_process_text_file', 'error', file_path, str(e))
        return lang.get('error_process_text_file').format(str(e)), None

# 7. handle_pdf_file: Processes PDF files
def handle_pdf_file(file_path):
    """Processes PDF files."""
    try:
        doc = fitz.open(file_path)
        text = [page.get_text("text") for page in doc]
        doc.close()
        log_message('pdf_file_processed', 'info', file_path)
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        log_message('error_process_pdf_file', 'error', file_path, str(e))
        return lang.get('error_process_pdf_file').format(file_path, str(e)), None

# 8. handle_word_file: Processes Word files
def handle_word_file(file_path):
    """Processes Word files."""
    try:
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        log_message('word_file_processed', 'info', file_path)
        return remove_headers_footers(text), file_path
    except Exception as e:
        log_message('error_process_word_file', 'error', file_path, str(e))
        return lang.get('error_process_word_file').format(file_path, str(e)), None

# 9. handle_ppt_file: Processes PowerPoint files
def handle_ppt_file(file_path):
    """Processes PowerPoint files."""
    try:
        ppt = Presentation(file_path)
        text = [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")]
        log_message('ppt_file_processed', 'info', file_path)
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        log_message('error_process_ppt_file', 'error', file_path, str(e))
        return lang.get('error_process_ppt_file').format(file_path, str(e)), None

# 10. handle_excel_file: Processes Excel files
def handle_excel_file(file_path):
    """Processes Excel files."""
    try:
        df = pd.read_excel(file_path)
        log_message('excel_file_processed', 'info', file_path)
        return df.to_csv(index=False), file_path
    except Exception as e:
        log_message('error_process_excel_file', 'error', file_path, str(e))
        return lang.get('error_process_excel_file').format(file_path, str(e)), None

# 11. handle_csv_file: Processes CSV files
def handle_csv_file(file_path):
    """Processes CSV files."""
    try:
        with open(file_path, mode='r', encoding='utf-8', newline='') as f:
            reader = csv.reader(f)
            data = list(reader)
        log_message('csv_file_processed', 'info', file_path)
        return '\n'.join([','.join(row) for row in data]), file_path
    except Exception as e:
        log_message('error_process_csv_file', 'error', file_path, str(e))
        return lang.get('error_process_csv_file').format(file_path, str(e)), None

# 12. handle_epub_file: Processes EPUB files
def handle_epub_file(file_path):
    """Processes EPUB files."""
    try:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        log_message('epub_file_processed', 'info', file_path)
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        log_message('error_process_epub_file', 'error', file_path, str(e))
        return lang.get('error_process_epub_file').format(file_path, str(e)), None

# 13. handle_xml_file: Processes XML files
def handle_xml_file(file_path):
    """Processes XML files."""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        texts = [elem.text for elem in root.iter() if elem.text is not None]
        log_message('xml_file_processed', 'info', file_path)
        return remove_headers_footers('\n'.join(texts)), file_path
    except Exception as e:
        log_message('error_process_xml_file', 'error', file_path, str(e))
        return lang.get('error_process_xml_file').format(file_path, str(e)), None

# 14. handle_audio_file: Processes audio files
def handle_audio_file(file_path):
    """Processes audio files."""
    if file_path.lower().endswith('.m4a'):
        sound = AudioSegment.from_file(file_path, format='m4a')
        wav_path = file_path.replace('.m4a', '.wav')
        sound.export(wav_path, format='wav')
        file_path = wav_path

    recognizer = sr.Recognizer()
    with sr.AudioFile(file_path) as source:
        audio_data = recognizer.record(source)
        try:
            text = recognizer.recognize_google(audio_data, language='it-IT')
            log_message('audio_file_processed', 'info', file_path)
            return text, file_path
        except sr.UnknownValueError:
            log_message('error_speech_not_understood', 'error', file_path)
            return lang.get('error_speech_not_understood').format(file_path), file_path
        except sr.RequestError as e:
            log_message('error_speech_recognition', 'error', file_path, str(e))
            return lang.get('error_speech_recognition').format(file_path, str(e)), file_path

# 15. handle_video_file: Processes video files
def handle_video_file(file_path):
    """Processes video files."""
    try:
        audio_path = extract_audio_from_video(file_path)
        text = transcribe_audio(audio_path)
        os.remove(audio_path)
        log_message('video_file_processed', 'info', file_path)
        return text, file_path
    except Exception as e:
        log_message('error_process_video_file', 'error', file_path, str(e))
        return lang.get('error_process_video_file').format(file_path, str(e)), None

# 16. extract_audio_from_video: Extracts audio from video
def extract_audio_from_video(video_path):
    """Extracts audio from video."""
    video = VideoFileClip(video_path)
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    audio_path = os.path.join(temp_dir, f"temp_audio_{timestamp}.wav")
    video.audio.write_audiofile(audio_path)
    return audio_path

# 17. transcribe_audio: Transcribes audio using Google Speech Recognition
def transcribe_audio(audio_path, language='it-IT'):
    """Transcribes audio using Google Speech Recognition."""
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio_data = recognizer.record(source)
        try:
            return recognizer.recognize_google(audio_data, language=language)
        except sr.UnknownValueError:
            return lang.get('error_speech_not_understood').format(audio_path)
        except sr.RequestError as e:
            return lang.get('error_speech_recognition').format(audio_path, str(e))

# 18. write_to_output: Writes content to output directory
def write_to_output(content, output_dir, file_index, original_path):
    """Writes content to output directory."""
    output_file_path = os.path.join(output_dir, f'model_{file_index}.txt')
    with open(output_file_path, 'a', encoding='utf-8') as file:
        file.write(f"\nOriginal file path: {original_path}\nFile content:\n{content}\n")
    log_message('output_written', 'info', output_file_path)
    return file_index + 1

# 19. handle_zip_file: Processes files within a ZIP archive
def handle_zip_file(zip_path):
    """Processes files within a ZIP archive."""
    try:
        with zipfile.ZipFile(zip_path, 'r') as z:
            z.extractall(temp_dir)
            extracted_files = z.namelist()
            for file_name in extracted_files:
                internal_path = os.path.join(temp_dir, file_name)
                if os.path.isfile(internal_path):
                    content, _ = handle_file(internal_path)
                    if content and not content.startswith("Unsupported"):
                        log_message('zip_file_processed', 'info', zip_path)
                        return f"{content} (from {file_name} in {zip_path})", zip_path
                    os.remove(internal_path)
        log_message('no_supported_files_found', 'warning', zip_path)
        return lang.get('no_supported_files_found').format(zip_path), None
    except Exception as e:
        log_message('error_process_file', 'error', zip_path, str(e))
        return lang.get('error_process_file').format(zip_path, str(e)), None

# 20. handle_file: Processes various file types
def handle_file(file_path):
    """Processes various file types."""
    extension = os.path.splitext(file_path)[1].lower()
    handler = {
        '.txt': handle_text_file,
        '.htm': handle_text_file,
        '.html': handle_text_file,
        '.srt': handle_text_file,
        '.pdf': handle_pdf_file,
        '.docx': handle_word_file,
        '.doc': handle_word_file,
        '.pptx': handle_ppt_file,
        '.ppt': handle_ppt_file,
        '.xls': handle_excel_file,
        '.xlsx': handle_excel_file,
        '.xml': handle_xml_file,
        '.gan': handle_xml_file,
        '.xsd': handle_xml_file,
        '.wav': handle_audio_file,
        '.mp3': handle_audio_file,
        '.m4a': handle_audio_file,
        '.mp4': handle_video_file,
        '.avi': handle_video_file,
        '.mov': handle_video_file,
        '.mkv': handle_video_file,
        '.mpeg': handle_video_file,
        '.mpg': handle_video_file,
        '.3gp': handle_video_file,
        '.csv': handle_csv_file,
        '.epub': handle_epub_file,
        '.zip': handle_zip_file
    }.get(extension)
    if handler:
        return handler(file_path)
    return lang.get('error_unknown_file_format').format(file_path), None

# 21. handle_directory: Processes all files in a directory
def handle_directory(directory_path, output_dir):
    """Processes all files in a directory."""
    file_index = 1
    for root, _, files in os.walk(directory_path):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            content, original_path = handle_file(file_path)
            if content and not content.startswith("Unsupported"):
                file_index = write_to_output(content, output_dir, file_index, original_path)

# 22. limit_files_search: Limits the search of files based on specific criteria
def limit_files_search(files, limit_search):
    """Limits the search of files based on specific criteria."""
    if limit_search == 'noLimit':
        return files
    if limit_search == 'lastProducedPerType':
        file_types = {}
        for file in files:
            file_type = os.path.splitext(file)[1]
            if (file_type not in file_types) or (os.path.getmtime(file) > os.path.getmtime(file_types[file_type])):
                file_types[file_type] = file
        return list(file_types.values())
    elif limit_search == 'lastProducedInFolder':
        if files:
            return [max(files, key=os.path.getmtime)]
    elif limit_search == 'lastProducedSimilarTitle':
        similar_titles = {}
        for file in files:
            base_name = os.path.splitext(file)[0]
            if base_name not in similar_titles or os.path.getmtime(file) > os.path.getmtime(similar_titles[base_name]):
                similarity = SequenceMatcher(None, base_name, os.path.splitext(similar_titles.get(base_name, ""))[0]).ratio()
                if similarity > 0.9:
                    similar_titles[base_name] = file
        return list(similar_titles.values())
    return files

# 23. process_text_with_keywords: Processes text with keywords and creates JSON data
def process_text_with_keywords(text, keywords):
    """Processes text with keywords and creates JSON data."""
    json_data = []
    keyword_positions = []

    # Find all keyword positions in text
    for keyword in keywords:
        pattern = re.compile(keyword, re.IGNORECASE)
        matches = list(pattern.finditer(text))
        for match in matches:
            keyword_positions.append((match.start(), match.end(), match.group()))

    # Sort keyword positions by their position in text
    keyword_positions.sort()

    # Add contents between keywords to JSON
    for i in range(len(keyword_positions)):
        start, end, matched_keyword = keyword_positions[i]
        next_start = keyword_positions[i + 1][0] if i + 1 < len(keyword_positions) else len(text)

        content = text[end:next_start].strip()
        json_data.append({"title": matched_keyword, "content": content})

    log_message('json_data_created', 'info')
    return json_data

# 24. write_json: Writes JSON data to file
def write_json(data, output_file):
    """Writes JSON data to file."""
    try:
        with open(output_file, 'w', encoding='utf-8') as json_file:
            json.dump(data, json_file, indent=4, ensure_ascii=False)
        log_message('json_file_written', 'info', output_file)
    except PermissionError:
        log_message('error_permission_denied', 'error', output_file)
    except Exception as e:
        log_message('error_write_json', 'error', output_file, str(e))

# 25. download_youtube_video: Downloads video from YouTube
def download_youtube_video(url, download_audio_only=False):
    """Downloads video from YouTube."""
    if not url.strip():
        log_message('error_download_url_empty', 'error')
        return None
    try:
        yt = YouTube(url)
        title = ''.join([c for c in yt.title if c.isalpha() or c.isdigit() or c == ' ']).rstrip()
        if download_audio_only:
            stream = yt.streams.filter(only_audio=True).first()
            file_path = stream.download(output_path=temp_dir, filename=f"{title}.mp3")
        else:
            stream = yt.streams.get_highest_resolution()
            file_path = stream.download(output_path=temp_dir, filename=f"{title}.mp4")
        log_message('success_download_youtube', 'info', file_path)
        return file_path
    except Exception as e:
        log_message('error_download_youtube', 'error', str(e))
        return None

# 26. download_vimeo_video: Downloads video from Vimeo
def download_vimeo_video(url):
    """Downloads video from Vimeo."""
    if not url.strip():
        log_message('error_download_url_empty', 'error')
        return None
    try:
        response = requests.get(url, stream=True)
        title = url.split("/")[-1]
        title = ''.join([c for c in title if c.isalpha() or c.isdigit() or c == ' ']).rstrip()
        file_path = os.path.join(temp_dir, f"{title}.mp4")
        with open(file_path, 'wb') as out_file:
            shutil.copyfileobj(response.raw, out_file)
        log_message('success_download_vimeo', 'info', file_path)
        return file_path
    except Exception as e:
        log_message('error_download_vimeo', 'error', str(e))
        return None

# 27. extract_audio: Extracts audio from video
def extract_audio(video_file):
    """Extracts audio from video."""
    try:
        video = VideoFileClip(video_file)
        audio = video.audio
        audio_file = os.path.join(temp_dir, "temp_audio.wav")
        audio.write_audiofile(audio_file, codec='pcm_s16le')
        video.close()
        log_message('audio_file_extracted', 'info', audio_file)
        return audio_file
    except Exception as e:
        log_message('error_extract_audio', 'error', str(e))
        return None

# 28. generate_srt: Generates SRT file from video
def generate_srt(video_file, output_file, language='it-IT'):
    """Generates SRT file from video."""
    try:
        # Extract audio from the video file
        audio_file = extract_audio(video_file)
        if not audio_file:
            raise FileNotFoundError(lang.get('error_extract_audio').format("Audio extraction failed."))

        sound = AudioSegment.from_wav(audio_file)
        chunks = split_on_silence(sound, min_silence_len=500, silence_thresh=sound.dBFS-14, keep_silence=500)
        
        with open(output_file, 'w') as file:
            start = 0
            for i, chunk in enumerate(chunks):
                chunk_filename = os.path.join(temp_dir, f"chunk_{i}.wav")
                chunk.export(chunk_filename, format="wav")
                with sr.AudioFile(chunk_filename) as source:
                    audio = sr.Recognizer().record(source)
                try:
                    text = sr.Recognizer().recognize_google(audio, language=language)
                    duration = len(chunk) / 1000
                    start_time = start
                    end_time = start + duration
                    file.write(f"{i+1}\n")
                    file.write(f"{format_time(start_time)} --> {format_time(end_time)}\n")
                    file.write(f"{text}\n\n")
                    start += duration
                    log_message('info_generated_srt_segment', 'info', i + 1)
                except sr.UnknownValueError:
                    log_message('warning_audio_not_understood', 'warning', i + 1)
                except sr.RequestError as e:
                    log_message('error_service_srt', 'error', i + 1, str(e))
                finally:
                    os.remove(chunk_filename)
        os.remove(audio_file)
    except Exception as e:
        log_message('error_generate_srt', 'error', str(e))

# 29. format_time: Formats time in SRT format
def format_time(seconds):
    """Formats time in SRT format."""
    hours, seconds = divmod(seconds, 3600)
    minutes, seconds = divmod(seconds, 60)
    milliseconds = int((seconds - int(seconds)) * 1000)
    return f"{int(hours):02}:{int(minutes):02}:{int(seconds):02},{milliseconds:03}"

# 30. upload_to_gdrive: Uploads a file to Google Drive
def upload_to_gdrive(file_path, folder_id, service):
    """Uploads a file to Google Drive."""
    file_metadata = {'name': os.path.basename(file_path), 'parents': [folder_id]}
    media = MediaFileUpload(file_path, resumable=True)
    file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    log_message('success_upload_gdrive', 'info', file.get('id'))

# 31. create_folder_on_gdrive: Creates a folder on Google Drive
def create_folder_on_gdrive(folder_name, parent_id, service):
    """Creates a folder on Google Drive."""
    folder_metadata = {'name': folder_name, 'mimeType': 'application/vnd.google-apps.folder', 'parents': [parent_id]}
    folder = service.files().create(body=folder_metadata, fields='id').execute()
    return folder.get('id')

# 32. upload_json_to_gdrive: Uploads JSON data to Google Drive
def upload_json_to_gdrive(json_data, file_name, folder_id, service):
    """Uploads JSON data to Google Drive."""
    json_path = os.path.join(temp_dir, file_name)
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(json_data, json_file, indent=4, ensure_ascii=False)
    upload_to_gdrive(json_path, folder_id, service)
    os.remove(json_path)

# 33. download_files_from_folder: Downloads all files from a Google Drive folder
def download_files_from_folder(folder_id, service, output_dir):
    """Downloads all files from a Google Drive folder."""
    results = service.files().list(q=f"'{folder_id}' in parents", spaces='drive', fields='files(id, name)').execute()
    items = results.get('files', [])
    for item in items:
        file_id = item['id']
        file_name = item['name']
        request = service.files().get_media(fileId=file_id)
        file_path = os.path.join(output_dir, file_name)
        with open(file_path, 'wb') as file:
            downloader = MediaIoBaseDownload(file, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
                log_message('download_progress', 'info', int(status.progress() * 100))
        log_message('success_download_gdrive', 'info', file_id, file_path)

# 34. download_from_gdrive: Downloads a file from Google Drive
def download_from_gdrive(file_id, service, output_path):
    """Downloads a file from Google Drive."""
    request = service.files().get_media(fileId=file_id)
    with open(output_path, 'wb') as file:
        downloader = MediaIoBaseDownload(file, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            log_message('download_progress', 'info', int(status.progress() * 100))
    log_message('success_download_gdrive', 'info', file_id, output_path)

# 35. download_all_files_from_gdrive: Downloads all files from a Google Drive folder and its subfolders
def download_all_files_from_gdrive(folder_id, service, output_dir):
    """Downloads all files from a Google Drive folder and its subfolders."""
    results = service.files().list(q=f"'{folder_id}' in parents", spaces='drive', fields='files(id, name, mimeType)').execute()
    items = results.get('files', [])
    for item in items:
        file_id = item['id']
        file_name = item['name']
        mime_type = item['mimeType']
        if mime_type == 'application/vnd.google-apps.folder':
            subfolder_path = os.path.join(output_dir, file_name)
            if not os.path.exists(subfolder_path):
                os.makedirs(subfolder_path)
            download_all_files_from_gdrive(file_id, service, subfolder_path)
        else:
            download_from_gdrive(file_id, service, os.path.join(output_dir, file_name))

# 36. upload_to_s3: Uploads a file to S3
def upload_to_s3(file_path, bucket_name, s3_client):
    """Uploads a file to S3."""
    s3_client.upload_file(file_path, bucket_name, os.path.basename(file_path))
    log_message('success_upload_s3', 'info', file_path, bucket_name)

# 37. download_from_s3: Downloads a file from S3
def download_from_s3(file_key, bucket_name, output_path, s3_client):
    """Downloads a file from S3."""
    s3_client.download_file(bucket_name, file_key, output_path)
    log_message('success_download_s3', 'info', file_key, bucket_name, output_path)

# 38. create_folder_on_s3: Creates a folder on S3
def create_folder_on_s3(folder_name, bucket_name, s3_client):
    """Creates a folder on S3."""
    s3_client.put_object(Bucket=bucket_name, Key=(folder_name + '/'))
    log_message('success_create_s3_folder', 'info', folder_name, bucket_name)

# 39. upload_json_to_s3: Uploads JSON data to S3
def upload_json_to_s3(json_data, file_name, bucket_name, s3_client):
    """Uploads JSON data to S3."""
    json_path = os.path.join(temp_dir, file_name)
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(json_data, json_file, indent=4, ensure_ascii=False)
    upload_to_s3(json_path, bucket_name, s3_client)
    os.remove(json_path)

# 40. download_directory_from_s3: Downloads a directory from S3
def download_directory_from_s3(bucket_name, s3_client, output_dir):
    """Downloads a directory from S3."""
    paginator = s3_client.get_paginator('list_objects_v2')
    for page in paginator.paginate(Bucket=bucket_name, Prefix=''):
        for obj in page.get('Contents', []):
            file_key = obj['Key']
            file_path = os.path.join(output_dir, file_key)
            if not os.path.exists(os.path.dirname(file_path)):
                os.makedirs(os.path.dirname(file_path))
            s3_client.download_file(bucket_name, file_key, file_path)
            log_message('success_download_s3', 'info', file_key, bucket_name, file_path)

# 41. upload_to_azure: Uploads a file to Azure Blob Storage
def upload_to_azure(file_path, container_name, blob_service_client):
    """Uploads a file to Azure Blob Storage."""
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=os.path.basename(file_path))
    with open(file_path, "rb") as data:
        blob_client.upload_blob(data)
    log_message('success_upload_azure', 'info', file_path, container_name)

# 42. download_from_azure: Downloads a file from Azure Blob Storage
def download_from_azure(blob_name, container_name, output_path, blob_service_client):
    """Downloads a file from Azure Blob Storage."""
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    with open(output_path, "wb") as file:
        data = blob_client.download_blob()
        data.readinto(file)
    log_message('success_download_azure', 'info', blob_name, container_name, output_path)

# 43. create_folder_on_azure: Creates a folder on Azure Blob Storage
def create_folder_on_azure(folder_name, container_name, blob_service_client):
    """Creates a folder on Azure Blob Storage."""
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=folder_name + '/')
    blob_client.upload_blob('', overwrite=True)
    log_message('success_create_azure_folder', 'info', folder_name, container_name)

# 44. upload_json_to_azure: Uploads JSON data to Azure Blob Storage
def upload_json_to_azure(json_data, file_name, container_name, blob_service_client):
    """Uploads JSON data to Azure Blob Storage."""
    json_path = os.path.join(temp_dir, file_name)
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(json_data, json_file, indent=4, ensure_ascii=False)
    upload_to_azure(json_path, container_name, blob_service_client)
    os.remove(json_path)

# 45. upload_to_aruba: Uploads a file to Aruba Cloud Object Storage
def upload_to_aruba(file_path, bucket_name, aruba_client):
    """Uploads a file to Aruba Cloud Object Storage."""
    aruba_client.upload_file(file_path, bucket_name, os.path.basename(file_path))
    log_message('success_upload_aruba', 'info', file_path, bucket_name)

# 46. download_from_aruba: Downloads a file from Aruba Cloud Object Storage
def download_from_aruba(file_key, bucket_name, output_path, aruba_client):
    """Downloads a file from Aruba Cloud Object Storage."""
    aruba_client.download_file(bucket_name, file_key, output_path)
    log_message('success_download_aruba', 'info', file_key, bucket_name, output_path)

# 47. create_folder_on_aruba: Creates a folder on Aruba Cloud Object Storage
def create_folder_on_aruba(folder_name, bucket_name, aruba_client):
    """Creates a folder on Aruba Cloud Object Storage."""
    aruba_client.put_object(Bucket=bucket_name, Key=(folder_name + '/'))
    log_message('success_create_aruba_folder', 'info', folder_name, bucket_name)

# 48. upload_json_to_aruba: Uploads JSON data to Aruba Cloud Object Storage
def upload_json_to_aruba(json_data, file_name, bucket_name, aruba_client):
    """Uploads JSON data to Aruba Cloud Object Storage."""
    json_path = os.path.join(temp_dir, file_name)
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(json_data, json_file, indent=4, ensure_ascii=False)
    upload_to_aruba(json_path, bucket_name, aruba_client)
    os.remove(json_path)

# 49. play_video_with_srt: Plays video with SRT subtitles
def play_video_with_srt(video_path, srt_path):
    """Plays video with SRT subtitles."""
    player = vlc.MediaPlayer()
    media = vlc.Media(video_path)
    media.add_option(f":sub-file={srt_path}")
    player.set_media(media)
    player.play()
    log_message('play_video_with_srt', 'info', video_path, srt_path)

# 50. play_video_from_command_line: Plays video from command line with VLC
def play_video_from_command_line(video_path):
    """Plays video from command line with VLC."""
    subprocess.call(['vlc', video_path])
    log_message('play_video_from_command_line', 'info', video_path)

# 51. download_directory_from_azure: Downloads a directory from Azure Blob Storage
def download_directory_from_azure(container_name, blob_service_client, output_dir):
    """Downloads a directory from Azure Blob Storage."""
    container_client = blob_service_client.get_container_client(container_name)
    blobs = container_client.list_blobs()
    for blob in blobs:
        file_path = os.path.join(output_dir, blob.name)
        if not os.path.exists(os.path.dirname(file_path)):
            os.makedirs(os.path.dirname(file_path))
        download_from_azure(blob.name, container_name, file_path, blob_service_client)

# 52. download_directory_from_aruba: Downloads a directory from Aruba Cloud Object Storage
def download_directory_from_aruba(bucket_name, aruba_client, output_dir):
    """Downloads a directory from Aruba Cloud Object Storage."""
    paginator = aruba_client.get_paginator('list_objects_v2')
    for page in paginator.paginate(Bucket=bucket_name, Prefix=''):
        for obj in page.get('Contents', []):
            file_key = obj['Key']
            file_path = os.path.join(output_dir, file_key)
            if not os.path.exists(os.path.dirname(file_path)):
                os.makedirs(os.path.dirname(file_path))
            aruba_client.download_file(bucket_name, file_key, file_path)
            log_message('success_download_aruba', 'info', file_key, bucket_name, file_path)

# 53. read_file_from_gdrive: Reads a file from Google Drive
def read_file_from_gdrive(file_id, service):
    """Reads a file from Google Drive."""
    request = service.files().get_media(fileId=file_id)
    data = request.execute()
    return data

# 54. read_file_from_s3: Reads a file from S3
def read_file_from_s3(file_key, bucket_name, s3_client):
    """Reads a file from S3."""
    obj = s3_client.get_object(Bucket=bucket_name, Key=file_key)
    data = obj['Body'].read()
    return data

# 55. read_file_from_azure: Reads a file from Azure Blob Storage
def read_file_from_azure(blob_name, container_name, blob_service_client):
    """Reads a file from Azure Blob Storage."""
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    data = blob_client.download_blob().readall()
    return data

# 56. read_file_from_aruba: Reads a file from Aruba Cloud Object Storage
def read_file_from_aruba(file_key, bucket_name, aruba_client):
    """Reads a file from Aruba Cloud Object Storage."""
    obj = aruba_client.get_object(Bucket=bucket_name, Key=file_key)
    data = obj['Body'].read()
    return data

# 57. create_container_if_not_exists: Creates an Azure container if it does not exist
def create_container_if_not_exists(container_name, blob_service_client):
    """Creates an Azure container if it does not exist."""
    container_client = blob_service_client.get_container_client(container_name)
    if not container_client.exists():
        container_client.create_container()
    log_message('success_create_azure_container', 'info', container_name)

# 58. upload_directory_to_azure: Uploads a directory to Azure Blob Storage
def upload_directory_to_azure(directory_path, container_name, blob_service_client):
    """Uploads a directory to Azure Blob Storage."""
    for root, _, files in os.walk(directory_path):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            blob_name = os.path.relpath(file_path, directory_path)
            upload_to_azure(file_path, container_name, blob_service_client)

# 59. launch_gui: Launches the GUI for the application
def launch_gui():
    """Launches the GUI for the application."""
    from PyQt5 import QtWidgets
    app = QtWidgets.QApplication([])
    window = QtWidgets.QMainWindow()
    window.setWindowTitle('Galora GUI')
    window.show()
    app.exec_()

# 60. extract_audio: Extracts audio from video (duplicata)
def extract_audio(video_file):
    """Extracts audio from video."""
    try:
        video = VideoFileClip(video_file)
        audio = video.audio
        audio_file = os.path.join(temp_dir, "temp_audio.wav")
        audio.write_audiofile(audio_file, codec='pcm_s16le')
        video.close()
        log_message('audio_file_extracted', 'info', audio_file)
        return audio_file
    except Exception as e:
        log_message('error_extract_audio', 'error', str(e))
        return None

# 61. process_video: Processes video files
def process_video(video_path, output_path):
    """Processes video files."""
    audio_path = extract_audio(video_path)
    srt_path = os.path.splitext(output_path)[0] + ".srt"
    generate_srt(audio_path, srt_path)
    play_video_with_srt(video_path, srt_path)

# 62. main: Main function to parse arguments and initiate processing
def main():
    print("Starting main function...")  # Stampa di debug
    parser = argparse.ArgumentParser(description="CLI Tool")
    parser.add_argument("--language", type=str, default="eng", help="Language code for localization")  # Funzione 4: load_translations
    parser.add_argument("--operation", type=str, help="Operation to perform")  # Funzioni varie, specificate sotto
    parser.add_argument("--file_path", type=str, help="Path to the file")  # Funzioni 6-15, 25-28
    parser.add_argument("--directory_path", type=str, help="Path to the directory")  # Funzioni 21, 22, 40, 51-53, 58
    parser.add_argument("--bucket_name", type=str, help="Bucket name for cloud storage")  # Funzioni 36-39, 45-48, 52
    parser.add_argument("--folder_id", type=str, help="Folder ID for Google Drive")  # Funzioni 30-35
    parser.add_argument("--file_id", type=str, help="File ID for Google Drive")  # Funzione 34
    parser.add_argument("--blob_name", type=str, help="Blob name for Azure Blob storage")  # Funzioni 42, 44
    parser.add_argument("--file_key", type=str, help="File key for S3/Aruba")  # Funzioni 37, 46
    parser.add_argument("--download_path", type=str, help="Path to download the file")  # Funzioni 34, 37, 42, 46
    parser.add_argument("--output_dir", type=str, help="Output directory")  # Funzioni 18, 21, 23, 58
    parser.add_argument("--keywords", nargs='*', help="Keywords for processing text files")  # Funzione 23
    parser.add_argument("--limit_search", type=str, default="noLimit", help="Limit file search criteria")  # Funzione 22
    parser.add_argument("--download_audio_only", action='store_true', help="Download audio only from video")  # Funzione 25
    parser.add_argument("--gui", action='store_true', help="Launch GUI interface")  # Funzione 59
    parser.add_argument("--play_video", action='store_true', help="Play video with SRT")  # Funzioni 49, 50
    parser.add_argument("--video_path", type=str, help="Path to the video file")  # Funzioni 15, 27, 49, 50
    parser.add_argument("--srt_path", type=str, help="Path to the SRT file")  # Funzione 49
    parser.add_argument("--file_name", type=str, help="Name of the file to read")  # Funzioni 53-56
    parser.add_argument("--upload_directory_to_azure", action="store_true", help="Upload directory to Azure Blob Storage")  # Funzione 58
    parser.add_argument("--download_directory_from_azure", action="store_true", help="Download directory from Azure Blob Storage")  # Funzione 51
    parser.add_argument("--container_name", type=str, help="Container name for Azure Blob storage")  # Funzioni 41, 43, 44, 51, 58
    parser.add_argument("--azure_directory", type=str, help="Directory path in Azure Blob storage")  # Funzione 51
    parser.add_argument("--url", type=str, help="URL of the video to download")  # Funzioni 25, 26
    parser.add_argument("--transcription_lang", type=str, default="en", help="Language for transcription")  # Funzioni 14, 17

    args = parser.parse_args()
    
    global translations
    translations = load_translations(args.language)  # Funzione 4
    print("Translations loaded.")  # Stampa di debug

    global config
    config = load_config("cli_tool")  # Funzione 3
    print("Config loaded.")  # Stampa di debug

    # Set the GOOGLE_APPLICATION_CREDENTIALS environment variable
    if config and 'google_application_credentials' in config:  # Funzioni 32-35
        os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = config['google_application_credentials']
        log_message('Google application credentials set: {}', 'info', config['google_application_credentials'])  # Funzione 2
        print("Google application credentials set.")  # Stampa di debug

    credentials = service_account.Credentials.from_service_account_file(  # Funzioni 32-35
        config['google_application_credentials'],
        scopes=config['gdrive_scopes']
    )
    print("Credentials loaded.")  # Stampa di debug

    configure_logger("cli_tool")  # Funzione 1
    print("Logger configured in main.")  # Stampa di debug
    # Launch GUI if --gui argument is passed
    if args.gui:
        launch_gui()  # Funzione 59
        return

    if args.operation == "process_video":  # Funzione 61
        process_video(
            url=args.url,
            file_path=args.file_path,
            download_audio_only=args.download_audio_only,
            transcription_lang=args.transcription_lang
        )
    else:
        print(f"Unknown operation: {args.operation}")
        logging.error(f"Unknown operation: {args.operation}")

    # Handle the play video operation
    if args.play_video:  # Funzioni 49, 50
        if args.video_path and args.srt_path:
            play_video_from_command_line(args.video_path, args.srt_path)  # Funzione 50
        else:
            log_message('Video or SRT path missing', 'error')  # Funzione 2
        return

    # Check that --operation is provided if --play_video is not specified
    if not args.operation and not args.upload_directory_to_azure and not args.download_directory_from_azure:
        parser.error('--operation is required unless --play_video is specified or --upload_directory_to_azure or --download_directory_from_azure is used')

    print(f"Operation: {args.operation}")  # Stampa di debug
    print(f"Upload to Azure: {args.upload_directory_to_azure}")  # Stampa di debug
    print(f"Download from Azure: {args.download_directory_from_azure}")  # Stampa di debug

    if args.upload_directory_to_azure:  # Funzione 58
        if config.get('use_azure', False):
            print(f"Uploading directory {args.directory_path} to Azure container {args.container_name}")  # Stampa di debug
            upload_directory_to_azure(args.directory_path, args.container_name)  # Funzione 58
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.download_directory_from_azure:  # Funzione 51
        if config.get('use_azure', False):
            print(f"Downloading directory {args.azure_directory} from Azure container {args.container_name} to {args.download_path}")  # Stampa di debug
            download_directory_from_azure(args.container_name, args.azure_directory, args.download_path)  # Funzione 51
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_gdrive":  # Funzione 30
        if config.get('use_gdrive', False):
            upload_to_gdrive(credentials, args.file_path, args.folder_id)  # Funzione 30
        else:
            log_message('Google Drive integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_gdrive":  # Funzione 34
        if config.get('use_gdrive', False):
            download_from_gdrive(credentials, args.file_id, args.download_path)  # Funzione 34
        else:
            log_message('Google Drive integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_all_gdrive":  # Funzione 35
        if config.get('use_gdrive', False):
            download_all_files_from_gdrive(credentials, args.output_dir)  # Funzione 35
        else:
            log_message('Google Drive integration is disabled', 'error')  # Funzione 2
    elif args.operation == "create_gdrive_folder":  # Funzione 31
        if config.get('use_gdrive', False):
            folder_id = create_folder_on_gdrive(credentials, args.folder_id)  # Funzione 31
            if folder_id:
                print(f"{folder_id}")
            else:
                log_message('Failed to create folder on Google Drive', 'error')  # Funzione 2
        else:
            log_message('Google Drive integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_json_to_gdrive":  # Funzione 32
        if config.get('use_gdrive', False):
            upload_json_to_gdrive(credentials, args.directory_path, args.folder_id)  # Funzione 32
        else:
            log_message('Google Drive integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_s3":  # Funzione 36
        if config.get('use_s3', False):
            upload_to_s3(args.file_path, args.bucket_name)  # Funzione 36
        else:
            log_message('S3 integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_s3":  # Funzione 37
        if config.get('use_s3', False):
            download_from_s3(args.file_key, args.bucket_name, args.download_path)  # Funzione 37
        else:
            log_message('S3 integration is disabled', 'error')  # Funzione 2
    elif args.operation == "create_s3_folder":  # Funzione 38
        if config.get('use_s3', False):
            create_folder_on_s3(args.bucket_name, args.folder_id)  # Funzione 38
        else:
            log_message('S3 integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_json_to_s3":  # Funzione 39
        if config.get('use_s3', False):
            upload_json_to_s3(args.directory_path, args.bucket_name, args.folder_id)  # Funzione 39
        else:
            log_message('S3 integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_azure":  # Funzione 41
        if config.get('use_azure', False):
            upload_to_azure(args.file_path, args.container_name)  # Funzione 41
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_azure":  # Funzione 42
        if config.get('use_azure', False):
            download_from_azure(args.blob_name, args.container_name, args.download_path)  # Funzione 42
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "create_azure_folder":  # Funzione 43
        if config.get('use_azure', False):
            create_folder_on_azure(args.container_name, args.folder_id)  # Funzione 43
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_json_to_azure":  # Funzione 44
        if config.get('use_azure', False):
            upload_json_to_azure(args.directory_path, args.container_name, args.folder_id)  # Funzione 44
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_aruba":  # Funzione 45
        if config.get('use_aruba', False):
            upload_to_aruba(args.file_path, args.bucket_name)  # Funzione 45
        else:
            log_message('Aruba integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_aruba":  # Funzione 46
        if config.get('use_aruba', False):
            download_from_aruba(args.file_key, args.bucket_name, args.download_path)  # Funzione 46
        else:
            log_message('Aruba integration is disabled', 'error')  # Funzione 2
    elif args.operation == "create_aruba_folder":  # Funzione 47
        if config.get('use_aruba', False):
            create_folder_on_aruba(args.bucket_name, args.folder_id)  # Funzione 47
        else:
            log_message('Aruba integration is disabled', 'error')  # Funzione 2
    elif args.operation == "upload_json_to_aruba":  # Funzione 48
        if config.get('use_aruba', False):
            upload_json_to_aruba(args.directory_path, args.bucket_name, args.folder_id)  # Funzione 48
        else:
            log_message('Aruba integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_youtube":  # Funzione 25
        download_youtube_video(args.file_path, args.download_audio_only)  # Funzione 25
    elif args.operation == "download_vimeo":  # Funzione 26
        download_vimeo_video(args.file_path)  # Funzione 26
    elif args.operation == "generate_srt":  # Funzione 28
        generate_srt(args.file_path, args.output_dir)  # Funzione 28
    elif args.operation == "handle_directory":  # Funzione 21
        handle_directory(args.directory_path, args.output_dir)  # Funzione 21
    elif args.operation == "download_s3_directory":  # Funzione 40
        if config.get('use_s3', False):
            download_directory_from_s3(args.bucket_name, args.directory_path, args.download_path)  # Funzione 40
        else:
            log_message('S3 integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_azure_directory":  # Funzione 51
        if config.get('use_azure', False):
            download_directory_from_azure(args.container_name, args.directory_path, args.download_path)  # Funzione 51
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "download_aruba_directory":  # Funzione 52
        if config.get('use_aruba', False):
            download_directory_from_aruba(args.bucket_name, args.directory_path, args.download_path)  # Funzione 52
        else:
            log_message('Aruba integration is disabled', 'error')  # Funzione 2
    elif args.operation == "read_gdrive_file":  # Funzione 53
        if config.get('use_gdrive', False):
            file_content = read_file_from_gdrive(credentials, args.folder_id, args.file_name)  # Funzione 53
            if file_content:
                with open(args.download_path, 'wb') as f:
                    f.write(file_content)
        else:
            log_message('Google Drive integration is disabled', 'error')  # Funzione 2
    elif args.operation == "read_s3_file":  # Funzione 54
        if config.get('use_s3', False):
            file_content = read_file_from_s3(args.bucket_name, args.directory_path, args.file_name)  # Funzione 54
            if file_content:
                with open(args.download_path, 'wb') as f:
                    f.write(file_content)
        else:
            log_message('S3 integration is disabled', 'error')  # Funzione 2
    elif args.operation == "read_azure_file":  # Funzione 55
        if config.get('use_azure', False):
            file_content = read_file_from_azure(args.container_name, args.directory_path, args.file_name)  # Funzione 55
            if file_content:
                with open(args.download_path, 'wb') as f:
                    f.write(file_content)
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "read_aruba_file":  # Funzione 56
        if config.get('use_aruba', False):
            file_content = read_file_from_aruba(args.bucket_name, args.directory_path, args.file_name)  # Funzione 56
            if file_content:
                with open(args.download_path, 'wb') as f:
                    f.write(file_content)
        else:
            log_message('Aruba integration is disabled', 'error')  # Funzione 2
    elif args.upload_directory_to_azure:  # Funzione 58
        if config.get('use_azure', False):
            print(f"Uploading directory {args.directory_path} to Azure container {args.container_name}")  # Stampa di debug
            upload_directory_to_azure(args.directory_path, args.container_name)  # Funzione 58
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.download_directory_from_azure:  # Funzione 51
        if config.get('use_azure', False):
            print(f"Downloading directory {args.azure_directory} from Azure container {args.container_name} to {args.download_path}")  # Stampa di debug
            download_directory_from_azure(args.container_name, args.azure_directory, args.download_path)  # Funzione 51
        else:
            log_message('Azure integration is disabled', 'error')  # Funzione 2
    elif args.operation == "process_keywords":  # Funzioni 21, 23
        if args.directory_path and args.output_dir and args.keywords:
            for root, _, files in os.walk(args.directory_path):
                for file_name in files:
                    file_path = os.path.join(root, file_name)
                    content, _ = handle_file(file_path)  # Funzione 20
                    if content and not content.startswith("Unsupported"):
                        json_data = process_text_with_keywords(content, args.keywords)  # Funzione 23
                        output_file = os.path.join(args.output_dir, f'{os.path.splitext(file_name)[0]}.json')
                        write_json(json_data, output_file)  # Funzione 24
        else:
            parser.error('--directory_path, --output_dir, and --keywords are required for process_keywords operation')
    else:
        log_message('Unknown operation: {}', 'error', args.operation)  # Funzione 2

    logging.shutdown()  # Assicurarsi che i log vengano scritti nel file
    print("Logging shutdown.")  # Stampa di debug

if __name__ == "__main__":
    main()
